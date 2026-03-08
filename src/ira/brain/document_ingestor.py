"""Document ingestion pipeline for Ira's knowledge base.

Walks the ``data/imports/`` directory tree, reads files in every supported
format (PDF, XLSX, DOCX, CSV, TXT), splits them into token-counted
overlapping chunks, and upserts the resulting :class:`KnowledgeItem` objects
into Qdrant via :class:`QdrantManager`.

A lightweight SQLite ledger (``data/ingested_files.db``) tracks which files
have already been processed so that re-running ingestion is idempotent.
"""

from __future__ import annotations

import asyncio
import csv
import hashlib
import io
import logging
import re
import sqlite3
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

import tiktoken

from ira.brain.qdrant_manager import QdrantManager
from ira.data.models import KnowledgeItem
from ira.exceptions import DatabaseError, IngestionError, LLMError, PathTraversalError

from typing import TYPE_CHECKING

if TYPE_CHECKING:
    from ira.brain.knowledge_graph import KnowledgeGraph

logger = logging.getLogger(__name__)

_SUPPORTED_EXTENSIONS = {".pdf", ".xlsx", ".docx", ".txt", ".csv", ".pptx", ".html", ".md"}
_DEFAULT_CHUNK_SIZE = 512
_DEFAULT_OVERLAP = 128
_TIKTOKEN_ENCODING = "cl100k_base"

_LEDGER_PATH = Path("data/ingested_files.db")

_CATEGORY_PATTERN = re.compile(r"^\d{2}_(.+)$")


# ── SQLite ledger ────────────────────────────────────────────────────────────


def _init_ledger(db_path: Path) -> sqlite3.Connection:
    db_path.parent.mkdir(parents=True, exist_ok=True)
    conn = sqlite3.connect(str(db_path))
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS ingested_files (
            path        TEXT PRIMARY KEY,
            hash        TEXT NOT NULL,
            chunk_count INTEGER NOT NULL,
            ingested_at TEXT NOT NULL
        )
        """
    )
    conn.commit()
    return conn


def _file_hash(path: Path) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for block in iter(lambda: f.read(1 << 16), b""):
            h.update(block)
    return h.hexdigest()


# ── file readers ─────────────────────────────────────────────────────────────


_MIN_USEFUL_CHARS = 50


def read_pdf(path: Path) -> str:
    text = _read_pdf_pypdf(path)
    if len(text.strip()) >= _MIN_USEFUL_CHARS:
        return text
    ocr_text = _read_pdf_document_ai(path)
    if ocr_text:
        return ocr_text
    return text


def _read_pdf_pypdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _read_pdf_document_ai(path: Path) -> str:
    """OCR fallback via Document AI for scanned/image-heavy PDFs."""
    try:
        from ira.systems.document_ai import DocumentAIService
        from ira.config import get_settings

        settings = get_settings()
        if not settings.document_ai.processor_id:
            return ""

        import asyncio
        import concurrent.futures

        svc = DocumentAIService()

        async def _ocr() -> str:
            await svc.connect()
            return await svc.extract_text(path.read_bytes())

        loop = None
        try:
            loop = asyncio.get_running_loop()
        except RuntimeError:
            pass

        if loop and loop.is_running():
            future = asyncio.run_coroutine_threadsafe(_ocr(), loop)
            try:
                return future.result(timeout=120)
            except (concurrent.futures.TimeoutError, Exception):
                logger.warning("Document AI OCR timed out for %s", path)
                return ""
        else:
            return asyncio.run(_ocr())
    except Exception:
        logger.warning("Document AI OCR fallback failed for %s", path, exc_info=True)
        return ""


def read_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    try:
        parts: list[str] = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            parts.append(f"[Sheet: {sheet}]")
            for row in ws.iter_rows(values_only=True):
                parts.append("\t".join(str(cell) if cell is not None else "" for cell in row))
        return "\n".join(parts)
    finally:
        wb.close()


def read_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    return "\n".join(para.text for para in doc.paragraphs)


def read_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def read_csv(path: Path) -> str:
    buf = io.StringIO()
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        for row in reader:
            buf.write("\t".join(row) + "\n")
    return buf.getvalue()


def read_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            parts.append(f"[Slide {i}]")
            parts.extend(texts)
    return "\n".join(parts)


def read_xls(path: Path) -> str:
    import xlrd

    wb = xlrd.open_workbook(str(path))
    parts: list[str] = []
    for sheet in wb.sheets():
        parts.append(f"[Sheet: {sheet.name}]")
        for row_idx in range(sheet.nrows):
            cells = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)]
            parts.append("\t".join(cells))
    return "\n".join(parts)


def _read_with_docling(path: Path) -> str:
    """Parse any supported document via Docling for high-fidelity extraction.

    Handles tables, reading order, formulas, and complex layouts far better
    than the legacy per-format readers.  Falls back to legacy readers on error.
    """
    try:
        from docling.document_converter import DocumentConverter

        converter = DocumentConverter()
        result = converter.convert(str(path))
        md = result.document.export_to_markdown()
        if md and len(md.strip()) >= _MIN_USEFUL_CHARS:
            return md
    except Exception:
        logger.debug("Docling failed for %s — falling back to legacy reader", path, exc_info=True)
    return ""


_DOCLING_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".html", ".md"}

_LEGACY_READERS = {
    ".pdf": read_pdf,
    ".xlsx": read_xlsx,
    ".xls": read_xls,
    ".docx": read_docx,
    ".txt": read_txt,
    ".csv": read_csv,
    ".pptx": read_pptx,
}


def _get_reader(ext: str):
    """Return a reader function: Docling-first for supported formats, else legacy."""
    if ext in _DOCLING_EXTENSIONS:
        def _docling_then_legacy(path: Path) -> str:
            text = _read_with_docling(path)
            if text:
                return text
            legacy = _LEGACY_READERS.get(ext)
            return legacy(path) if legacy else ""
        return _docling_then_legacy
    return _LEGACY_READERS.get(ext)


_READERS = {ext: _get_reader(ext) or fn for ext, fn in _LEGACY_READERS.items()}
_READERS.update({ext: _get_reader(ext) for ext in _DOCLING_EXTENSIONS if ext not in _READERS})


# ── chunking ─────────────────────────────────────────────────────────────────


def _chunk_text_tiktoken(
    text: str,
    chunk_size: int = _DEFAULT_CHUNK_SIZE,
    overlap: int = _DEFAULT_OVERLAP,
) -> list[str]:
    """Legacy fixed-size token chunking with tiktoken.  Used as fallback."""
    enc = tiktoken.get_encoding(_TIKTOKEN_ENCODING)
    tokens = enc.encode(text)

    if len(tokens) <= chunk_size:
        return [text]

    chunks: list[str] = []
    start = 0
    while start < len(tokens):
        end = min(start + chunk_size, len(tokens))
        chunks.append(enc.decode(tokens[start:end]))
        if end == len(tokens):
            break
        start += chunk_size - overlap

    return chunks


_TABLE_PATTERN = re.compile(
    r"((?:^\|.+\|$\n?){2,})",
    re.MULTILINE,
)
_TABLE_SENTINEL_START = "\n\n<!-- TABLE_BOUNDARY -->\n"
_TABLE_SENTINEL_END = "\n<!-- /TABLE_BOUNDARY -->\n\n"


def _protect_tables(text: str) -> tuple[str, list[str]]:
    """Wrap Markdown tables in sentinels so the chunker won't split them."""
    tables: list[str] = []

    def _replace(m: re.Match) -> str:
        tables.append(m.group(1))
        idx = len(tables) - 1
        return f"{_TABLE_SENTINEL_START}__TABLE_{idx}__{_TABLE_SENTINEL_END}"

    return _TABLE_PATTERN.sub(_replace, text), tables


def _restore_tables(chunks: list[str], tables: list[str]) -> list[str]:
    """Replace table placeholders with original table text."""
    restored: list[str] = []
    for chunk in chunks:
        for idx, table in enumerate(tables):
            placeholder = f"__TABLE_{idx}__"
            if placeholder in chunk:
                chunk = chunk.replace(
                    f"{_TABLE_SENTINEL_START}{placeholder}{_TABLE_SENTINEL_END}",
                    f"\n\n{table}\n\n",
                )
        restored.append(chunk.strip())
    return [c for c in restored if c]


def _get_voyage_embeddings():
    """Build a Chonkie VoyageAIEmbeddings instance using our configured key."""
    try:
        from chonkie.embeddings import VoyageAIEmbeddings
        from ira.config import get_settings

        settings = get_settings()
        api_key = settings.embedding.api_key.get_secret_value()
        model = settings.embedding.model
        if not api_key:
            return None
        return VoyageAIEmbeddings(model=model, api_key=api_key)
    except Exception:
        logger.debug("VoyageAI embeddings for Chonkie unavailable", exc_info=True)
        return None


def chunk_text(
    text: str,
    chunk_size: int = _DEFAULT_CHUNK_SIZE,
    overlap: int = _DEFAULT_OVERLAP,
) -> list[str]:
    """Split *text* into semantically coherent chunks via Chonkie.

    Uses SemanticChunker with VoyageAI embeddings (matching our retrieval
    embedding space) when available.  Markdown tables are protected from
    being split across chunk boundaries.  Falls back to the legacy
    fixed-size tiktoken chunker on error.
    """
    protected_text, tables = _protect_tables(text)

    try:
        from chonkie import SemanticChunker

        embedding_model = _get_voyage_embeddings() or "minishlab/potion-base-32M"
        chunker = SemanticChunker(
            embedding_model=embedding_model,
            chunk_size=chunk_size,
            threshold=0.7,
            similarity_window=3,
        )
        chunks = chunker.chunk(protected_text)
        result = [c.text for c in chunks if c.text.strip()]
        if result:
            return _restore_tables(result, tables)
    except Exception:
        logger.debug("Chonkie semantic chunking failed — using tiktoken fallback", exc_info=True)

    fallback = _chunk_text_tiktoken(protected_text, chunk_size, overlap)
    return _restore_tables(fallback, tables)


# ── category extraction ──────────────────────────────────────────────────────


def _category_from_path(file_path: Path, base_path: Path) -> str:
    """Derive a source_category slug from the first directory under *base_path*.

    Expects folder names like ``01_Quotes_and_Proposals``.  Returns the part
    after the numeric prefix, lowercased (e.g. ``quotes_and_proposals``).
    Falls back to ``"uncategorised"`` if the pattern doesn't match.
    """
    try:
        relative = file_path.relative_to(base_path)
        top_dir = relative.parts[0] if relative.parts else ""
    except ValueError:
        return "uncategorised"

    m = _CATEGORY_PATTERN.match(top_dir)
    return m.group(1).lower() if m else top_dir.lower() or "uncategorised"


# ── ingestor class ───────────────────────────────────────────────────────────


class DocumentIngestor:
    """Reads, chunks, and upserts documents from the imports directory."""

    def __init__(
        self,
        qdrant: QdrantManager,
        *,
        knowledge_graph: "KnowledgeGraph | None" = None,
        ledger_path: Path = _LEDGER_PATH,
    ) -> None:
        self._qdrant = qdrant
        self._graph = knowledge_graph
        self._ledger = _init_ledger(ledger_path)

    # ── discovery ────────────────────────────────────────────────────────

    def discover_files(self, base_path: str = "data/imports") -> list[dict[str, Any]]:
        """Walk *base_path* and return metadata for every supported file."""
        base = Path(base_path).resolve()
        project_root = Path(__file__).resolve().parents[3]
        data_root = project_root / "data"
        if not base.is_relative_to(data_root):
            raise PathTraversalError(f"Path {base_path} is outside the data directory")

        root = Path(base_path)
        if not root.exists():
            logger.warning("Import directory does not exist: %s", root)
            return []

        files: list[dict[str, Any]] = []
        for p in sorted(root.rglob("*")):
            if not p.is_file():
                continue
            ext = p.suffix.lower()
            if ext not in _SUPPORTED_EXTENSIONS:
                continue
            files.append(
                {
                    "path": str(p),
                    "category": _category_from_path(p, root),
                    "extension": ext,
                    "size": p.stat().st_size,
                }
            )

        logger.info("Discovered %d importable files under %s", len(files), root)
        return files

    # ── contextual retrieval ─────────────────────────────────────────────

    _CONTEXT_PREVIEW_CHARS = 3000

    async def _generate_document_context(self, text: str, filename: str) -> str:
        """Generate a 2-sentence document summary for contextual retrieval.

        Prepended to every chunk so isolated chunks retain the global
        context of their source document (Anthropic's Contextual Retrieval).
        """
        try:
            from ira.prompt_loader import load_prompt
            from ira.services.llm_client import get_llm_client

            system = load_prompt("document_context")
            preview = text[: self._CONTEXT_PREVIEW_CHARS]
            user_msg = f"Filename: {filename}\n\n{preview}"
            summary = await get_llm_client().generate_text(
                system, user_msg,
                temperature=0.0, max_tokens=200,
                name="document_context",
            )
            return summary.strip()
        except Exception:
            logger.warning(
                "Document context generation failed for %s — proceeding without context",
                filename, exc_info=True,
            )
            return ""

    # ── single-file ingestion ────────────────────────────────────────────

    def is_already_ingested(self, file_info: dict[str, Any]) -> bool:
        """Check whether a file has already been ingested with the same hash."""
        path = Path(file_info["path"])
        return self._already_ingested(str(path), _file_hash(path))

    async def ingest_file(
        self,
        file_info: dict[str, Any],
        *,
        force: bool = False,
        chunk_size: int = _DEFAULT_CHUNK_SIZE,
        overlap: int = _DEFAULT_OVERLAP,
    ) -> int:
        """Read, chunk, and upsert one file.  Returns the chunk count.

        When *force* is ``True`` the ledger check is skipped and any
        existing Qdrant points for this source path are deleted first.
        """
        path = Path(file_info["path"])
        category = file_info["category"]
        ext = file_info["extension"]

        current_hash = await asyncio.to_thread(_file_hash, path)
        if not force and self._already_ingested(str(path), current_hash):
            logger.debug("Skipping already-ingested file: %s", path)
            return 0

        if force:
            await self._qdrant.delete_by_source(str(path))

        reader = _get_reader(ext) or _READERS.get(ext)
        if reader is None:
            logger.warning("No reader for extension '%s': %s", ext, path)
            return 0

        text = await asyncio.to_thread(reader, path)

        if ext == ".pdf" and len(text.strip()) < _MIN_USEFUL_CHARS:
            ocr_text = await self._ocr_fallback(path)
            if ocr_text:
                text = ocr_text

        if not text.strip():
            logger.warning("Empty content after reading: %s", path)
            return 0

        doc_context = await self._generate_document_context(text, path.name)

        chunks = chunk_text(text, chunk_size=chunk_size, overlap=overlap)

        items = [
            KnowledgeItem(
                source=str(path),
                source_category=category,
                content=(
                    f"[Source Context: {doc_context}]\n\n{chunk}"
                    if doc_context else chunk
                ),
                metadata={
                    "chunk_index": i,
                    "total_chunks": len(chunks),
                    "extension": ext,
                    "document_context": doc_context or "",
                },
            )
            for i, chunk in enumerate(chunks)
        ]

        upserted = await self._qdrant.upsert_items(items)

        if self._graph is not None:
            await self._extract_and_store_entities(text, str(path))

        self._record_ingestion(str(path), current_hash, upserted)
        logger.info("Ingested %s -> %d chunks (category: %s)", path, upserted, category)
        return upserted

    # ── bulk ingestion ───────────────────────────────────────────────────

    async def ingest_all(
        self,
        base_path: str = "data/imports",
        *,
        force: bool = False,
    ) -> dict[str, Any]:
        """Process every supported file under *base_path*.

        Returns a summary dict with keys ``files_processed``,
        ``files_skipped``, ``total_chunks``, ``per_category``, and
        ``errors`` (a list of ``{"path": ..., "error": ...}`` dicts).
        """
        files = self.discover_files(base_path)

        files_processed = 0
        files_skipped = 0
        total_chunks = 0
        per_category: dict[str, int] = {}
        errors: list[dict[str, str]] = []

        for file_info in files:
            try:
                n = await self.ingest_file(file_info, force=force)
                if n > 0:
                    files_processed += 1
                    total_chunks += n
                    cat = file_info["category"]
                    per_category[cat] = per_category.get(cat, 0) + n
                else:
                    files_skipped += 1
            except (IngestionError, Exception) as exc:
                logger.exception("Failed to ingest %s", file_info["path"])
                errors.append({"path": file_info["path"], "error": str(exc)})

        summary: dict[str, Any] = {
            "files_processed": files_processed,
            "files_skipped": files_skipped,
            "total_chunks": total_chunks,
            "per_category": per_category,
            "errors": errors,
        }
        logger.info("Ingestion complete: %s", summary)
        return summary

    # ── Document AI OCR fallback ─────────────────────────────────────────

    async def _ocr_fallback(self, path: Path) -> str:
        """Try Document AI OCR when pypdf yields insufficient text."""
        try:
            from ira.systems.document_ai import DocumentAIService

            svc = DocumentAIService()
            await svc.connect()
            if not svc.available:
                return ""
            text = await svc.extract_text(path.read_bytes())
            if text.strip():
                logger.info("Document AI OCR recovered text for %s (%d chars)", path, len(text))
            return text
        except Exception:
            logger.warning("Document AI OCR fallback failed for %s", path, exc_info=True)
            return ""

    async def reingest_scanned_pdfs(
        self,
        base_path: str = "data/imports",
        *,
        min_file_size: int = 5 * 1024 * 1024,
    ) -> dict[str, Any]:
        """Re-ingest PDFs that are likely scanned (large files with poor text).

        Finds PDFs over *min_file_size* bytes, checks if pypdf yields little
        text, and re-ingests them with Document AI OCR.
        """
        files = self.discover_files(base_path)
        candidates = [
            f for f in files
            if f["extension"] == ".pdf" and f["size"] >= min_file_size
        ]
        logger.info(
            "Found %d PDF candidates for OCR re-ingestion (>%d bytes)",
            len(candidates), min_file_size,
        )

        reingested = 0
        skipped = 0
        errors: list[dict[str, str]] = []

        for file_info in candidates:
            path = Path(file_info["path"])
            try:
                pypdf_text = await asyncio.to_thread(_read_pdf_pypdf, path)
                if len(pypdf_text.strip()) >= _MIN_USEFUL_CHARS:
                    skipped += 1
                    continue

                ocr_text = await self._ocr_fallback(path)
                if not ocr_text.strip():
                    skipped += 1
                    continue

                n = await self.ingest_file(file_info, force=True)
                if n > 0:
                    reingested += 1
                    logger.info("Re-ingested scanned PDF: %s -> %d chunks", path, n)
            except Exception as exc:
                logger.exception("Failed to re-ingest %s", path)
                errors.append({"path": str(path), "error": str(exc)})

        summary = {
            "candidates": len(candidates),
            "reingested": reingested,
            "skipped": skipped,
            "errors": errors,
        }
        logger.info("Scanned PDF re-ingestion complete: %s", summary)
        return summary

    # ── entity extraction ─────────────────────────────────────────────────

    async def _extract_and_store_entities(self, text: str, source: str) -> None:
        """Extract entities from document text and store them in Neo4j.

        Extraction priority: GraphRAG (schema-bound, with entity resolution)
        > legacy LLM extraction > GLiNER (fast local fallback).
        Results from all available extractors are merged and deduplicated.
        """
        assert self._graph is not None

        try:
            from ira.brain.entity_extractor import extract_entities_gliner
            gliner_entities = await asyncio.to_thread(extract_entities_gliner, text)
        except Exception:
            logger.debug("GLiNER extraction failed for %s — using LLM only", source)
            gliner_entities = {"companies": [], "people": [], "machines": [], "relationships": []}

        try:
            entities = await self._graph.extract_entities_from_text(text)
        except (LLMError, Exception):
            logger.warning("Entity extraction failed for %s — using GLiNER results only", source)
            entities = gliner_entities

        entities = self._merge_entity_results(gliner_entities, entities)

        for company in entities.get("companies", []):
            try:
                await self._graph.add_company(
                    name=company.get("name", ""),
                    region=company.get("region", ""),
                    industry=company.get("industry", ""),
                )
            except (DatabaseError, Exception):
                logger.warning("Failed to add company from %s: %s", source, company)

        for person in entities.get("people", []):
            try:
                await self._graph.add_person(
                    name=person.get("name", ""),
                    email=person.get("email", ""),
                    company_name=person.get("company", ""),
                    role=person.get("role", ""),
                )
            except (DatabaseError, Exception):
                logger.warning("Failed to add person from %s: %s", source, person)

        for machine in entities.get("machines", []):
            try:
                await self._graph.add_machine(
                    model=machine.get("model", ""),
                    category=machine.get("category", ""),
                    description=machine.get("description", ""),
                )
            except (DatabaseError, Exception):
                logger.warning("Failed to add machine from %s: %s", source, machine)

        rel_count = 0
        for rel in entities.get("relationships", []):
            try:
                ok = await self._graph.add_relationship(
                    from_type=rel.get("from_type", ""),
                    from_key=rel.get("from_key", ""),
                    rel_type=rel.get("rel", ""),
                    to_type=rel.get("to_type", ""),
                    to_key=rel.get("to_key", ""),
                    properties={k: v for k, v in rel.items()
                                if k not in ("from_type", "from_key", "rel", "to_type", "to_key")},
                )
                if ok:
                    rel_count += 1
            except (DatabaseError, Exception):
                logger.warning("Failed to add relationship from %s: %s", source, rel)

        logger.info(
            "Extracted entities from %s: %d companies, %d people, %d machines, %d relationships",
            source,
            len(entities.get("companies", [])),
            len(entities.get("people", [])),
            len(entities.get("machines", [])),
            rel_count,
        )

    @staticmethod
    def _merge_entity_results(a: dict[str, Any], b: dict[str, Any]) -> dict[str, Any]:
        """Merge two entity extraction results, deduplicating by key fields."""
        def _dedup(items: list[dict], key_field: str) -> list[dict]:
            seen: set[str] = set()
            out: list[dict] = []
            for item in items:
                k = item.get(key_field, "").lower().strip()
                if k and k not in seen:
                    seen.add(k)
                    out.append(item)
            return out

        return {
            "companies": _dedup(a.get("companies", []) + b.get("companies", []), "name"),
            "people": _dedup(a.get("people", []) + b.get("people", []), "name"),
            "machines": _dedup(a.get("machines", []) + b.get("machines", []), "model"),
            "relationships": a.get("relationships", []) + b.get("relationships", []),
        }

    # ── ledger helpers ───────────────────────────────────────────────────

    def _file_hash_for(self, file_info: dict[str, Any]) -> str:
        """Return the SHA-256 hash for a discovered file."""
        return _file_hash(Path(file_info["path"]))

    def _already_ingested(self, path: str, file_hash: str) -> bool:
        row = self._ledger.execute(
            "SELECT hash FROM ingested_files WHERE path = ?", (path,)
        ).fetchone()
        return row is not None and row[0] == file_hash

    def _record_ingestion(self, path: str, file_hash: str, chunk_count: int) -> None:
        self._ledger.execute(
            """
            INSERT INTO ingested_files (path, hash, chunk_count, ingested_at)
            VALUES (?, ?, ?, ?)
            ON CONFLICT(path) DO UPDATE SET
                hash = excluded.hash,
                chunk_count = excluded.chunk_count,
                ingested_at = excluded.ingested_at
            """,
            (path, file_hash, chunk_count, datetime.now(timezone.utc).isoformat()),
        )
        self._ledger.commit()

    def close(self) -> None:
        self._ledger.close()
